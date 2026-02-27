from __future__ import annotations

import json
import os
import re
import sqlite3
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from flask import Flask, jsonify, render_template, request

BASE_DIR = Path(__file__).resolve().parent
DB_PATH = BASE_DIR / "flashcards.db"
UPLOAD_DIR = BASE_DIR / "uploads"
ALLOWED_EXTENSIONS = {".pdf", ".pptx"}
STATUS_CHOICES = {"know", "kind_of_know", "dont_know"}
MAX_CARD_COUNT = 80

STOPWORDS = {
    "about",
    "after",
    "again",
    "against",
    "along",
    "among",
    "because",
    "before",
    "between",
    "could",
    "every",
    "first",
    "found",
    "important",
    "include",
    "including",
    "lecture",
    "material",
    "might",
    "other",
    "paper",
    "project",
    "result",
    "results",
    "should",
    "slides",
    "their",
    "there",
    "these",
    "those",
    "through",
    "using",
    "where",
    "which",
    "while",
    "would",
}


def create_app() -> Flask:
    app = Flask(__name__)
    app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024

    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    init_db()

    @app.get("/")
    def index() -> str:
        return render_template("index.html")

    @app.get("/api/health")
    def health() -> Any:
        return jsonify({"ok": True, "timestamp": utc_now()})

    @app.get("/api/projects")
    def get_projects() -> Any:
        with get_conn() as conn:
            rows = conn.execute(
                """
                SELECT
                    p.id,
                    p.name,
                    p.description,
                    p.created_at,
                    COUNT(f.id) AS total_cards,
                    COALESCE(SUM(CASE WHEN f.status = 'know' THEN 1 ELSE 0 END), 0) AS know_count,
                    COALESCE(SUM(CASE WHEN f.status = 'kind_of_know' THEN 1 ELSE 0 END), 0) AS kind_of_know_count,
                    COALESCE(SUM(CASE WHEN f.status = 'dont_know' THEN 1 ELSE 0 END), 0) AS dont_know_count
                FROM projects p
                LEFT JOIN flashcards f ON f.project_id = p.id
                GROUP BY p.id
                ORDER BY p.created_at DESC
                """
            ).fetchall()

        return jsonify({"projects": [dict(row) for row in rows]})

    @app.post("/api/projects")
    def create_project() -> Any:
        payload = request.get_json(silent=True) or {}
        name = str(payload.get("name", "")).strip()
        description = str(payload.get("description", "")).strip()

        if not name:
            return jsonify({"error": "Project name is required."}), 400

        with get_conn() as conn:
            cursor = conn.execute(
                "INSERT INTO projects (name, description, created_at) VALUES (?, ?, ?)",
                (name, description, utc_now()),
            )
            project_id = cursor.lastrowid
            conn.commit()

            project = conn.execute(
                "SELECT id, name, description, created_at FROM projects WHERE id = ?",
                (project_id,),
            ).fetchone()

        return jsonify({"project": dict(project)}), 201

    @app.get("/api/projects/<int:project_id>")
    def get_project(project_id: int) -> Any:
        with get_conn() as conn:
            project = conn.execute(
                "SELECT id, name, description, created_at FROM projects WHERE id = ?",
                (project_id,),
            ).fetchone()
            if project is None:
                return jsonify({"error": "Project not found."}), 404

            sources = conn.execute(
                """
                SELECT id, filename, extension, created_at
                FROM source_materials
                WHERE project_id = ?
                ORDER BY created_at DESC
                """,
                (project_id,),
            ).fetchall()

            counts = conn.execute(
                """
                SELECT
                    COUNT(*) AS total_cards,
                    COALESCE(SUM(CASE WHEN status = 'know' THEN 1 ELSE 0 END), 0) AS know_count,
                    COALESCE(SUM(CASE WHEN status = 'kind_of_know' THEN 1 ELSE 0 END), 0) AS kind_of_know_count,
                    COALESCE(SUM(CASE WHEN status = 'dont_know' THEN 1 ELSE 0 END), 0) AS dont_know_count
                FROM flashcards
                WHERE project_id = ?
                """,
                (project_id,),
            ).fetchone()

        return jsonify(
            {
                "project": dict(project),
                "sources": [dict(row) for row in sources],
                "stats": dict(counts),
            }
        )

    @app.get("/api/projects/<int:project_id>/cards")
    def get_project_cards(project_id: int) -> Any:
        status = request.args.get("status", "all").strip()

        if status != "all" and status not in STATUS_CHOICES:
            return jsonify({"error": "Invalid status filter."}), 400

        with get_conn() as conn:
            project = conn.execute(
                "SELECT id FROM projects WHERE id = ?",
                (project_id,),
            ).fetchone()
            if project is None:
                return jsonify({"error": "Project not found."}), 404

            if status == "all":
                rows = conn.execute(
                    """
                    SELECT id, project_id, source_id, question, answer, status, created_at, updated_at
                    FROM flashcards
                    WHERE project_id = ?
                    ORDER BY created_at ASC
                    """,
                    (project_id,),
                ).fetchall()
            else:
                rows = conn.execute(
                    """
                    SELECT id, project_id, source_id, question, answer, status, created_at, updated_at
                    FROM flashcards
                    WHERE project_id = ? AND status = ?
                    ORDER BY created_at ASC
                    """,
                    (project_id, status),
                ).fetchall()

        return jsonify({"cards": [dict(row) for row in rows]})

    @app.post("/api/projects/<int:project_id>/cards")
    def add_manual_card(project_id: int) -> Any:
        payload = request.get_json(silent=True) or {}
        question = str(payload.get("question", "")).strip()
        answer = str(payload.get("answer", "")).strip()

        if not question or not answer:
            return jsonify({"error": "Question and answer are required."}), 400

        with get_conn() as conn:
            project = conn.execute(
                "SELECT id FROM projects WHERE id = ?",
                (project_id,),
            ).fetchone()
            if project is None:
                return jsonify({"error": "Project not found."}), 404

            cursor = conn.execute(
                """
                INSERT INTO flashcards
                    (project_id, source_id, question, answer, status, created_at, updated_at)
                VALUES (?, NULL, ?, ?, 'dont_know', ?, ?)
                """,
                (project_id, question, answer, utc_now(), utc_now()),
            )
            conn.commit()

            card = conn.execute(
                """
                SELECT id, project_id, source_id, question, answer, status, created_at, updated_at
                FROM flashcards
                WHERE id = ?
                """,
                (cursor.lastrowid,),
            ).fetchone()

        return jsonify({"card": dict(card)}), 201

    @app.patch("/api/cards/<int:card_id>/status")
    def update_card_status(card_id: int) -> Any:
        payload = request.get_json(silent=True) or {}
        status = str(payload.get("status", "")).strip()

        if status not in STATUS_CHOICES:
            return jsonify({"error": "Invalid status value."}), 400

        with get_conn() as conn:
            existing = conn.execute(
                "SELECT id FROM flashcards WHERE id = ?",
                (card_id,),
            ).fetchone()
            if existing is None:
                return jsonify({"error": "Card not found."}), 404

            conn.execute(
                "UPDATE flashcards SET status = ?, updated_at = ? WHERE id = ?",
                (status, utc_now(), card_id),
            )
            conn.commit()

            card = conn.execute(
                """
                SELECT id, project_id, source_id, question, answer, status, created_at, updated_at
                FROM flashcards
                WHERE id = ?
                """,
                (card_id,),
            ).fetchone()

        return jsonify({"card": dict(card)})

    @app.post("/api/projects/<int:project_id>/upload")
    def upload_and_generate(project_id: int) -> Any:
        if "file" not in request.files:
            return jsonify({"error": "File is required."}), 400

        file = request.files["file"]
        if not file.filename:
            return jsonify({"error": "No file selected."}), 400

        extension = Path(file.filename).suffix.lower()
        if extension not in ALLOWED_EXTENSIONS:
            return (
                jsonify(
                    {
                        "error": "Unsupported file type. Use PDF or PPTX.",
                        "accepted": sorted(ALLOWED_EXTENSIONS),
                    }
                ),
                400,
            )

        card_count = clamp_card_count(request.form.get("card_count"))

        with get_conn() as conn:
            project = conn.execute(
                "SELECT id FROM projects WHERE id = ?",
                (project_id,),
            ).fetchone()
            if project is None:
                return jsonify({"error": "Project not found."}), 404

        safe_name = sanitize_filename(file.filename)
        storage_name = f"{uuid.uuid4().hex}_{safe_name}"
        saved_path = UPLOAD_DIR / storage_name
        file.save(saved_path)

        try:
            extracted_text = extract_text(saved_path, extension)
        except Exception as exc:  # noqa: BLE001
            saved_path.unlink(missing_ok=True)
            return jsonify({"error": f"Failed to read document: {exc}"}), 400

        if len(extracted_text.strip()) < 200:
            saved_path.unlink(missing_ok=True)
            return (
                jsonify(
                    {
                        "error": "Could not extract enough text to generate flashcards."
                    }
                ),
                400,
            )

        generated_cards, model_used = generate_flashcards(extracted_text, card_count)
        if not generated_cards:
            return jsonify({"error": "Could not generate flashcards."}), 500

        now = utc_now()
        with get_conn() as conn:
            cursor = conn.execute(
                """
                INSERT INTO source_materials
                    (project_id, filename, storage_name, extension, extracted_text, created_at)
                VALUES (?, ?, ?, ?, ?, ?)
                """,
                (
                    project_id,
                    safe_name,
                    storage_name,
                    extension,
                    extracted_text,
                    now,
                ),
            )
            source_id = cursor.lastrowid

            for card in generated_cards:
                conn.execute(
                    """
                    INSERT INTO flashcards
                        (project_id, source_id, question, answer, status, created_at, updated_at)
                    VALUES (?, ?, ?, ?, 'dont_know', ?, ?)
                    """,
                    (
                        project_id,
                        source_id,
                        card["question"].strip(),
                        card["answer"].strip(),
                        now,
                        now,
                    ),
                )
            conn.commit()

        return jsonify(
            {
                "source": {
                    "id": source_id,
                    "filename": safe_name,
                    "extension": extension,
                },
                "generated_count": len(generated_cards),
                "generator": model_used,
            }
        )

    return app


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def get_conn() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON;")
    return conn


def init_db() -> None:
    with get_conn() as conn:
        conn.executescript(
            """
            CREATE TABLE IF NOT EXISTS projects (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                name TEXT NOT NULL,
                description TEXT,
                created_at TEXT NOT NULL
            );

            CREATE TABLE IF NOT EXISTS source_materials (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                project_id INTEGER NOT NULL,
                filename TEXT NOT NULL,
                storage_name TEXT NOT NULL,
                extension TEXT NOT NULL,
                extracted_text TEXT NOT NULL,
                created_at TEXT NOT NULL,
                FOREIGN KEY(project_id) REFERENCES projects(id) ON DELETE CASCADE
            );

            CREATE TABLE IF NOT EXISTS flashcards (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                project_id INTEGER NOT NULL,
                source_id INTEGER,
                question TEXT NOT NULL,
                answer TEXT NOT NULL,
                status TEXT NOT NULL DEFAULT 'dont_know',
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                FOREIGN KEY(project_id) REFERENCES projects(id) ON DELETE CASCADE,
                FOREIGN KEY(source_id) REFERENCES source_materials(id) ON DELETE SET NULL
            );

            CREATE INDEX IF NOT EXISTS idx_flashcards_project_id ON flashcards(project_id);
            CREATE INDEX IF NOT EXISTS idx_flashcards_status ON flashcards(status);
            CREATE INDEX IF NOT EXISTS idx_source_project_id ON source_materials(project_id);
            """
        )
        conn.commit()


def sanitize_filename(filename: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9._-]", "_", filename)
    return cleaned[:180] or "upload_file"


def clamp_card_count(raw_value: Any) -> int:
    try:
        value = int(raw_value)
    except (TypeError, ValueError):
        return 20

    return max(5, min(value, MAX_CARD_COUNT))


def extract_text(path: Path, extension: str) -> str:
    if extension == ".pdf":
        return extract_pdf_text(path)
    if extension == ".pptx":
        return extract_pptx_text(path)

    raise ValueError("Unsupported file extension")


def extract_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("pypdf is not installed") from exc

    reader = PdfReader(str(path))
    chunks: list[str] = []
    for page in reader.pages:
        chunks.append(page.extract_text() or "")

    return normalize_text("\n".join(chunks))


def extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("python-pptx is not installed") from exc

    presentation = Presentation(str(path))
    chunks: list[str] = []

    for slide_idx, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            text = str(text).strip()
            if text:
                slide_lines.append(text)

        if slide_lines:
            chunks.append(f"Slide {slide_idx}: " + " | ".join(slide_lines))

    return normalize_text("\n".join(chunks))


def normalize_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def generate_flashcards(text: str, count: int) -> tuple[list[dict[str, str]], str]:
    api_key = os.getenv("OPENAI_API_KEY")
    if api_key:
        try:
            cards = generate_with_openai(text, count)
            if cards:
                return cards[:count], os.getenv("OPENAI_MODEL", "gpt-4.1-mini")
        except Exception:  # noqa: BLE001
            pass

    return generate_locally(text, count), "local-fallback"


def generate_with_openai(text: str, count: int) -> list[dict[str, str]]:
    try:
        from openai import OpenAI
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("openai package is not installed") from exc

    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    model = os.getenv("OPENAI_MODEL", "gpt-4.1-mini")

    lecture_text = text[:50000]
    prompt = (
        "Generate detailed, exam-grade flashcards from this lecture material. "
        "Make questions varied: conceptual, mechanisms, compare/contrast, and application. "
        "Return strict JSON with this shape only: "
        '{"cards": [{"question": "...", "answer": "..."}]}. '
        f"Generate exactly {count} cards."
    )

    response = client.responses.create(
        model=model,
        temperature=0.2,
        input=[
            {
                "role": "system",
                "content": "You create rigorous university and PhD-level study flashcards.",
            },
            {
                "role": "user",
                "content": f"{prompt}\n\nLECTURE MATERIAL:\n{lecture_text}",
            },
        ],
    )

    response_text = extract_response_text(response)
    cards = parse_cards_json(response_text)
    if not cards:
        raise ValueError("Model returned no valid cards")

    return cards[:count]


def extract_response_text(response: Any) -> str:
    text = getattr(response, "output_text", None)
    if text:
        return str(text)

    collected: list[str] = []
    for item in getattr(response, "output", []) or []:
        for content in getattr(item, "content", []) or []:
            maybe_text = getattr(content, "text", None)
            if maybe_text:
                collected.append(str(maybe_text))

    return "\n".join(collected)


def parse_cards_json(raw: str) -> list[dict[str, str]]:
    cleaned = raw.strip()
    if cleaned.startswith("```"):
        cleaned = cleaned.strip("`")
        cleaned = cleaned.replace("json", "", 1).strip()

    payload = json.loads(cleaned)
    rows = payload.get("cards", []) if isinstance(payload, dict) else []

    cards: list[dict[str, str]] = []
    for item in rows:
        if not isinstance(item, dict):
            continue
        question = str(item.get("question", "")).strip()
        answer = str(item.get("answer", "")).strip()
        if question and answer:
            cards.append({"question": question, "answer": answer})

    return cards


def generate_locally(text: str, count: int) -> list[dict[str, str]]:
    chunks = chunk_text(text)
    if not chunks:
        return []

    cards: list[dict[str, str]] = []
    seen_questions: set[str] = set()

    for chunk in chunks:
        question = f"Explain the core idea behind: {infer_title(chunk)}"
        answer = summarize_chunk(chunk, max_length=540)
        if question not in seen_questions and answer:
            cards.append({"question": question, "answer": answer})
            seen_questions.add(question)
        if len(cards) >= count:
            return cards[:count]

    for sentence in sentence_stream(text):
        for term in extract_terms(sentence):
            question = f"What is {term}, and why does it matter in this lecture context?"
            if question in seen_questions:
                continue
            answer = sentence.strip()
            if len(answer) < 40:
                continue
            cards.append({"question": question, "answer": answer})
            seen_questions.add(question)
            if len(cards) >= count:
                return cards[:count]

    while len(cards) < count:
        idx = len(cards) + 1
        cards.append(
            {
                "question": f"Key takeaway #{idx}: what should you remember?",
                "answer": summarize_chunk(text, max_length=500),
            }
        )

    return cards[:count]


def chunk_text(text: str, target_size: int = 800) -> list[str]:
    paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]
    if not paragraphs:
        return []

    chunks: list[str] = []
    buffer = ""

    for paragraph in paragraphs:
        if len(paragraph) < 50:
            continue
        candidate = f"{buffer}\n\n{paragraph}".strip() if buffer else paragraph
        if len(candidate) <= target_size:
            buffer = candidate
        else:
            if buffer:
                chunks.append(buffer)
            buffer = paragraph

    if buffer:
        chunks.append(buffer)

    return chunks[:MAX_CARD_COUNT]


def infer_title(chunk: str) -> str:
    first_line = chunk.split("\n", 1)[0]
    first_sentence = re.split(r"(?<=[.!?])\s+", first_line)[0]
    compact = re.sub(r"\s+", " ", first_sentence).strip(" :-")
    if len(compact) > 90:
        compact = compact[:87].rstrip() + "..."
    return compact or "this lecture section"


def summarize_chunk(chunk: str, max_length: int = 500) -> str:
    sentences = [s.strip() for s in sentence_stream(chunk) if len(s.strip()) > 20]
    if not sentences:
        return chunk[:max_length]

    summary = " ".join(sentences[:4]).strip()
    if len(summary) > max_length:
        summary = summary[: max_length - 3].rstrip() + "..."
    return summary


def sentence_stream(text: str) -> list[str]:
    return [s for s in re.split(r"(?<=[.!?])\s+", text) if s.strip()]


def extract_terms(sentence: str) -> list[str]:
    words = re.findall(r"\b[A-Za-z][A-Za-z0-9-]{5,}\b", sentence)
    terms: list[str] = []
    for word in words:
        normalized = word.lower()
        if normalized in STOPWORDS:
            continue
        if normalized.endswith("ing") and len(normalized) < 8:
            continue
        if word not in terms:
            terms.append(word)
        if len(terms) >= 2:
            break
    return terms


app = create_app()


if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=8000)
